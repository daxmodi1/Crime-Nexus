import os
import zipfile
import tempfile
import shutil
from typing import List, Tuple
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from core.models import EvidenceMetadata
from config.settings import settings

# Supported evidence file extensions
SUPPORTED_EXTENSIONS = [".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".csv", ".json", ".log", ".rtf"]


class SafeDocLoader:
    """
    Loader for old .doc files (Word 97-2003 format).
    On Windows: Uses pywin32 to leverage Microsoft Word COM
    On Linux: Uses antiword, catdoc, or LibreOffice
    """
    def __init__(self, file_path: str):
        self.file_path = os.path.abspath(file_path)

    def load(self):
        content = None
        
        # Method 1: Try pywin32 (Windows with Word installed - most reliable)
        try:
            import win32com.client
            import pythoncom
            
            pythoncom.CoInitialize()
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                
                doc = word.Documents.Open(self.file_path)
                content = doc.Content.Text
                doc.Close(False)
                word.Quit()
                
                if content and content.strip():
                    print(f"[DOC] Loaded via Word COM: {len(content)} chars")
                    return [Document(
                        page_content=content,
                        metadata={"source": self.file_path}
                    )]
            finally:
                pythoncom.CoUninitialize()
        except ImportError:
            print("[DOC] pywin32 not installed, trying alternative methods...")
        except Exception as e:
            print(f"[DOC] Word COM failed: {e}")
        
        # Method 2: Try antiword (Linux/Mac)
        try:
            import subprocess
            result = subprocess.run(
                ['antiword', self.file_path],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                content = result.stdout
                print(f"[DOC] Loaded via antiword: {len(content)} chars")
                return [Document(
                    page_content=content,
                    metadata={"source": self.file_path}
                )]
        except FileNotFoundError:
            print("[DOC] antiword not found in PATH")
        except Exception as e:
            print(f"[DOC] antiword failed: {e}")
        
        # Method 3: Try using catdoc (Linux/Mac)
        try:
            import subprocess
            result = subprocess.run(
                ['catdoc', self.file_path],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                content = result.stdout
                print(f"[DOC] Loaded via catdoc: {len(content)} chars")
                return [Document(
                    page_content=content,
                    metadata={"source": self.file_path}
                )]
        except FileNotFoundError:
            print("[DOC] catdoc not found in PATH")
        except Exception as e:
            print(f"[DOC] catdoc failed: {e}")
        
        # Method 4: Try using LibreOffice to convert
        try:
            import subprocess
            import tempfile
            
            with tempfile.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    ['soffice', '--headless', '--convert-to', 'txt:Text', '--outdir', tmpdir, self.file_path],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                
                # Find the output txt file
                txt_file = os.path.join(tmpdir, os.path.splitext(os.path.basename(self.file_path))[0] + '.txt')
                if os.path.exists(txt_file):
                    with open(txt_file, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    if content.strip():
                        print(f"[DOC] Loaded via LibreOffice: {len(content)} chars")
                        return [Document(
                            page_content=content,
                            metadata={"source": self.file_path}
                        )]
        except FileNotFoundError:
            print("[DOC] LibreOffice (soffice) not found in PATH")
        except Exception as e:
            print(f"[DOC] LibreOffice conversion failed: {e}")
        
        raise RuntimeError(
            f"Failed to extract text from .doc file: {self.file_path}. "
            "On Windows, make sure Microsoft Word is installed. "
            "On Linux/Mac, install antiword, catdoc, or LibreOffice"
        )


class SafeDocxLoader:
    """
    Custom robust loader for .docx files.
    Tries multiple methods to extract text.
    """
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self):
        content = None
        
        # Method 1: Try docx2txt
        try:
            import docx2txt
            content = docx2txt.process(self.file_path)
            if content and content.strip():
                print(f"[DOCX] Loaded via docx2txt: {len(content)} chars")
                return [Document(
                    page_content=content,
                    metadata={"source": self.file_path}
                )]
        except Exception as e:
            print(f"[DOCX] docx2txt failed: {e}")
        
        # Method 2: Try python-docx
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(self.file_path)
            paragraphs = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        paragraphs.append(" | ".join(row_text))
            
            content = "\n\n".join(paragraphs)
            if content and content.strip():
                print(f"[DOCX] Loaded via python-docx: {len(content)} chars")
                return [Document(
                    page_content=content,
                    metadata={"source": self.file_path}
                )]
        except Exception as e:
            print(f"[DOCX] python-docx failed: {e}")
        
        # Method 3: Extract raw XML text from DOCX (it's a ZIP file)
        try:
            import zipfile
            import re
            
            with zipfile.ZipFile(self.file_path, 'r') as z:
                xml_content = z.read('word/document.xml').decode('utf-8')
                # Extract text between XML tags
                text = re.sub(r'<[^>]+>', ' ', xml_content)
                text = re.sub(r'\s+', ' ', text).strip()
                
                if text:
                    print(f"[DOCX] Loaded via XML extraction: {len(text)} chars")
                    return [Document(
                        page_content=text,
                        metadata={"source": self.file_path}
                    )]
        except Exception as e:
            print(f"[DOCX] XML extraction failed: {e}")
        
        raise RuntimeError(f"Failed to extract text from DOCX: {self.file_path}")


class SafePptxLoader:
    """
    Custom lightweight loader for .pptx files.
    Requires: uv add python-pptx
    """
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self):
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx library not found. Please run: uv add python-pptx")

        try:
            prs = Presentation(self.file_path)
            full_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    full_text.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))

            return [Document(
                page_content="\n\n".join(full_text),
                metadata={"source": self.file_path}
            )]
            
        except Exception as e:
            raise RuntimeError(f"Failed to process PPTX: {str(e)}")


class SafePptLoader:
    """
    Loader for old .ppt files (PowerPoint 97-2003 format).
    Requires: LibreOffice or textract
    """
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self):
        content = None
        
        # Method 1: Try textract
        try:
            import textract
            content = textract.process(self.file_path).decode('utf-8')
            if content and content.strip():
                print(f"[PPT] Loaded via textract: {len(content)} chars")
                return [Document(
                    page_content=content,
                    metadata={"source": self.file_path}
                )]
        except ImportError:
            print("[PPT] textract not installed, trying LibreOffice...")
        except Exception as e:
            print(f"[PPT] textract failed: {e}")
        
        # Method 2: Convert via LibreOffice
        try:
            import subprocess
            import tempfile
            
            with tempfile.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    ['soffice', '--headless', '--convert-to', 'txt:Text', '--outdir', tmpdir, self.file_path],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                
                txt_file = os.path.join(tmpdir, os.path.splitext(os.path.basename(self.file_path))[0] + '.txt')
                if os.path.exists(txt_file):
                    with open(txt_file, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    if content.strip():
                        print(f"[PPT] Loaded via LibreOffice: {len(content)} chars")
                        return [Document(
                            page_content=content,
                            metadata={"source": self.file_path}
                        )]
        except FileNotFoundError:
            print("[PPT] LibreOffice (soffice) not found in PATH")
        except Exception as e:
            print(f"[PPT] LibreOffice conversion failed: {e}")
        
        raise RuntimeError(
            f"Failed to extract text from .ppt file: {self.file_path}. "
            "Please install LibreOffice or textract"
        )


class SafeRtfLoader:
    """
    Loader for RTF (Rich Text Format) files.
    Requires: striprtf or textract
    """
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self):
        content = None
        
        # Method 1: Try striprtf
        try:
            from striprtf.striprtf import rtf_to_text
            with open(self.file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            content = rtf_to_text(rtf_content)
            if content and content.strip():
                print(f"[RTF] Loaded via striprtf: {len(content)} chars")
                return [Document(
                    page_content=content,
                    metadata={"source": self.file_path}
                )]
        except ImportError:
            print("[RTF] striprtf not installed, trying alternatives...")
        except Exception as e:
            print(f"[RTF] striprtf failed: {e}")
        
        # Method 2: Try textract
        try:
            import textract
            content = textract.process(self.file_path).decode('utf-8')
            if content and content.strip():
                print(f"[RTF] Loaded via textract: {len(content)} chars")
                return [Document(
                    page_content=content,
                    metadata={"source": self.file_path}
                )]
        except ImportError:
            print("[RTF] textract not installed")
        except Exception as e:
            print(f"[RTF] textract failed: {e}")
        
        raise RuntimeError(
            f"Failed to extract text from .rtf file: {self.file_path}. "
            "Please install striprtf: uv add striprtf"
        )


class TextFileLoader:
    """
    Simple loader for text-based files (.txt, .csv, .json, .log)
    """
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self):
        try:
            with open(self.file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            return [Document(
                page_content=content,
                metadata={"source": self.file_path}
            )]
        except Exception as e:
            raise RuntimeError(f"Failed to read text file: {str(e)}")


def is_supported_file(filename: str) -> bool:
    """Check if a file has a supported extension."""
    ext = os.path.splitext(filename)[1].lower()
    return ext in SUPPORTED_EXTENSIONS


def extract_zip_file(zip_path: str, extract_to: str) -> List[str]:
    """
    Extract a ZIP file and return list of extracted file paths.
    Only extracts supported file types.
    """
    extracted_files = []
    
    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            # Get list of files in the zip
            for file_info in zip_ref.infolist():
                # Skip directories
                if file_info.is_dir():
                    continue
                
                # Get the filename (handle nested paths)
                filename = os.path.basename(file_info.filename)
                
                # Skip hidden files and unsupported types
                if filename.startswith('.') or filename.startswith('__'):
                    continue
                
                if not is_supported_file(filename):
                    print(f"Skipping unsupported file: {filename}")
                    continue
                
                # Extract file
                # Create a flat structure (all files in same directory)
                source = zip_ref.open(file_info)
                target_path = os.path.join(extract_to, filename)
                
                # Handle duplicate filenames
                base, ext = os.path.splitext(filename)
                counter = 1
                while os.path.exists(target_path):
                    target_path = os.path.join(extract_to, f"{base}_{counter}{ext}")
                    counter += 1
                
                with open(target_path, 'wb') as target:
                    shutil.copyfileobj(source, target)
                
                extracted_files.append(target_path)
                print(f"Extracted: {filename}")
        
        return extracted_files
    
    except zipfile.BadZipFile:
        raise ValueError("Invalid or corrupted ZIP file")
    except Exception as e:
        raise RuntimeError(f"Failed to extract ZIP file: {str(e)}")


def load_and_split_document(file_path: str, metadata: EvidenceMetadata):
    """
    Robust loader that handles PDF, DOCX, DOC, PPTX, PPT, and text files.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    loader = None
    
    try:
        # 1. Select the specific loader for the file type
        if ext == ".pdf":
            loader = PyPDFLoader(file_path)
            
        elif ext == ".docx":
            loader = SafeDocxLoader(file_path)
            
        elif ext == ".doc":
            loader = SafeDocLoader(file_path)
            
        elif ext == ".pptx":
            loader = SafePptxLoader(file_path)
            
        elif ext == ".ppt":
            # For old .ppt format, try converting via LibreOffice
            loader = SafePptLoader(file_path)
            
        elif ext == ".rtf":
            loader = SafeRtfLoader(file_path)
            
        elif ext in [".txt", ".csv", ".json", ".log"]:
            loader = TextFileLoader(file_path)
            
        else:
            raise ValueError(f"Unsupported forensic report format: {ext}")
            
        # 2. Load the content
        print(f"Loading {ext} file: {file_path}...")
        raw_docs = loader.load()
        
        if not raw_docs:
            print("Warning: Document was empty or could not be parsed.")
            return []
    
        # 3. Attach Forensic Metadata (Chain of Custody)
        for doc in raw_docs:
            doc.metadata.update(metadata.to_dict())
            
        # 4. Split into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.CHUNK_SIZE,
            chunk_overlap=settings.CHUNK_OVERLAP,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        chunks = text_splitter.split_documents(raw_docs)
        print(f"Successfully processed {len(chunks)} chunks from {os.path.basename(file_path)}")
        return chunks

    except ImportError as e:
        error_msg = f"Missing dependency for {ext}. Error: {str(e)}"
        print(error_msg)
        raise ImportError("Please run: uv add python-docx python-pptx docx2txt")
        
    except Exception as e:
        print(f"Critical Error loading {file_path}: {e}")
        raise e